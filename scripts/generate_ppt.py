"""
生成开题展示PPT - 多模态具身VLA Agent
12页，适合上台汇报
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ========== 配色方案 ==========
DARK_BG = RGBColor(0x1A, 0x1A, 0x2E)        # 深蓝黑背景
ACCENT_BLUE = RGBColor(0x00, 0x78, 0xD4)     # 亮蓝强调色
ACCENT_GREEN = RGBColor(0x10, 0x7C, 0x10)    # 绿色
ACCENT_ORANGE = RGBColor(0xD8, 0x7B, 0x00)   # 橙色
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
MEDIUM_GRAY = RGBColor(0x99, 0x99, 0x99)
CARD_BG = RGBColor(0x25, 0x25, 0x3A)         # 卡片背景

def set_slide_bg(slide, color):
    """设置幻灯片背景颜色"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape_bg(slide, left, top, width, height, color, transparency=0.3):
    """添加矩形背景作为卡片"""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=18, 
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="微软雅黑"):
    """添加文本框"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_bullet_list(slide, left, top, width, height, items, font_size=16, 
                    color=WHITE, bullet_color=ACCENT_BLUE, spacing=Pt(8)):
    """添加项目符号列表"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "微软雅黑"
        p.space_after = spacing
        p.level = 0
    return txBox

def add_table(slide, rows, cols, left, top, width, height):
    """添加表格"""
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    return table_shape.table

def style_table_cell(cell, text, font_size=14, color=WHITE, bold=False, bg_color=None, alignment=PP_ALIGN.CENTER):
    """设置单元格样式"""
    cell.text = ""
    p = cell.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "微软雅黑"
    p.alignment = alignment
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    if bg_color:
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg_color

# ============================================================
# 第1页：封面
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
set_slide_bg(slide, DARK_BG)

# 装饰线条
line1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(2.5), Inches(11.3), Inches(0.03))
line1.fill.solid()
line1.fill.fore_color.rgb = ACCENT_BLUE
line1.line.fill.background()

add_text_box(slide, Inches(1), Inches(1), Inches(11.3), Inches(1.5),
             "多模态大模型原理与应用 · 期中开题", font_size=20, 
             color=MEDIUM_GRAY, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(2.8), Inches(11.3), Inches(1.8),
             "基于视觉-语言-动作(VLA)的\n具身智能导航Agent", font_size=36, 
             color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(4.8), Inches(11.3), Inches(1),
             "在AI2-THOR仿真环境中实现多模态目标导航与交互", font_size=22, 
             color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)

line2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(5.8), Inches(11.3), Inches(0.03))
line2.fill.solid()
line2.fill.fore_color.rgb = ACCENT_BLUE
line2.line.fill.background()

add_text_box(slide, Inches(1), Inches(6.2), Inches(11.3), Inches(1),
             "Qwen3-VL 8B  ·  CLIP  ·  MobileSAM  ·  PPO", font_size=18, 
             color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# ============================================================
# 第2页：背景与动机
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(5), Inches(0.8),
             "背景与动机", font_size=32, color=ACCENT_BLUE, bold=True)

# 左侧：传统VLM的局限
card1 = add_shape_bg(slide, Inches(0.8), Inches(1.3), Inches(5.5), Inches(2.8), CARD_BG)
add_text_box(slide, Inches(1), Inches(1.4), Inches(5), Inches(0.5),
             "传统视觉-语言模型的局限", font_size=20, color=ACCENT_ORANGE, bold=True)
add_bullet_list(slide, Inches(1.2), Inches(2), Inches(5), Inches(2),
                ["\u274c \u505c\u7559\u5728\u201c\u770b\u548c\u8bf4\u201d\u5c42\u9762", 
                 "\u274c \u65e0\u6cd5\u5728\u7269\u7406\u73af\u5883\u4e2d\u6267\u884c\u52a8\u4f5c",
                 "\u274c \u7f3a\u4e4f\u957f\u7a0b\u51b3\u7b56\u89c4\u5212\u80fd\u529b", 
                 "\u274c \u4e0e\u73af\u5883\u7684\u4ea4\u4e92\u80fd\u529b\u7f3a\u5931"],
                font_size=16, color=LIGHT_GRAY)

# 右侧：具身AI的优势
card2 = add_shape_bg(slide, Inches(7), Inches(1.3), Inches(5.5), Inches(2.8), CARD_BG)
add_text_box(slide, Inches(7.2), Inches(1.4), Inches(5), Inches(0.5),
             "具身智能(AI)的突破", font_size=20, color=ACCENT_GREEN, bold=True)
add_bullet_list(slide, Inches(7.4), Inches(2), Inches(5), Inches(2),
                ["✅ 视觉+语言+动作的统一框架", "✅ 在仿真/真实环境中主动交互",
                 "✅ 从感知到决策的闭环", "✅ 家庭服务/仓储物流等实际应用"],
                font_size=16, color=LIGHT_GRAY)

# 底部：核心挑战
card3 = add_shape_bg(slide, Inches(0.8), Inches(4.5), Inches(11.7), Inches(2.5), CARD_BG)
add_text_box(slide, Inches(1), Inches(4.6), Inches(5), Inches(0.5),
             "核心技术挑战", font_size=20, color=WHITE, bold=True)
add_bullet_list(slide, Inches(1.2), Inches(5.2), Inches(11), Inches(1.8),
                ["1. 多模态特征对齐：视觉(图像)、语言(指令)、状态(位置)三种异构特征如何融合",
                 "2. 长程决策规划：导航需要数十步连续决策，存在信用分配(credit assignment)难题",
                 "3. 泛化能力：训练在有限场景，测试需泛化到未见过的房间布局和物体摆放"],
                font_size=16, color=LIGHT_GRAY, spacing=Pt(10))

# ============================================================
# 第3页：任务定义
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(5), Inches(0.8),
             "任务定义", font_size=32, color=ACCENT_BLUE, bold=True)

# 输入
card1 = add_shape_bg(slide, Inches(0.8), Inches(1.3), Inches(3.7), Inches(3), CARD_BG)
add_text_box(slide, Inches(1), Inches(1.4), Inches(3), Inches(0.5),
             "📥 输入", font_size=22, color=ACCENT_GREEN, bold=True)
add_bullet_list(slide, Inches(1.1), Inches(2), Inches(3.3), Inches(2.2),
                ["自然语言指令", '"Go to the microwave"', "第一人称RGB图像", 
                 "300×300像素", "Agent自身状态", "[x,y,z,rotation]"],
                font_size=15, color=LIGHT_GRAY, spacing=Pt(4))

# 输出
card2 = add_shape_bg(slide, Inches(4.8), Inches(1.3), Inches(3.7), Inches(3), CARD_BG)
add_text_box(slide, Inches(5), Inches(1.4), Inches(3), Inches(0.5),
             "📤 输出", font_size=22, color=ACCENT_ORANGE, bold=True)
add_bullet_list(slide, Inches(5.1), Inches(2), Inches(3.3), Inches(2.2),
                ["离散动作序列", "MoveAhead", "RotateLeft / RotateRight", 
                 "LookUp / LookDown", "Pickup (拾取)"],
                font_size=15, color=LIGHT_GRAY, spacing=Pt(4))

# 任务
card3 = add_shape_bg(slide, Inches(8.8), Inches(1.3), Inches(3.7), Inches(3), CARD_BG)
add_text_box(slide, Inches(9), Inches(1.4), Inches(3), Inches(0.5),
             "🎯 任务类型", font_size=22, color=ACCENT_BLUE, bold=True)
add_bullet_list(slide, Inches(9.1), Inches(2), Inches(3.3), Inches(2.2),
                ["ObjectNav 目标导航", "导航至目标物体可视范围", "距离≤1.5m且无障碍",
                 "Pickup 拾取交互", "到达目标后执行拾取"],
                font_size=15, color=LIGHT_GRAY, spacing=Pt(4))

# 成功判定
card4 = add_shape_bg(slide, Inches(0.8), Inches(4.6), Inches(11.7), Inches(2.4), CARD_BG)
add_text_box(slide, Inches(1), Inches(4.7), Inches(5), Inches(0.5),
             "✅ 成功判定标准", font_size=22, color=WHITE, bold=True)

table = add_table(slide, 3, 3, Inches(1.2), Inches(5.4), Inches(11), Inches(1.3))
table.columns[0].width = Inches(3)
table.columns[1].width = Inches(5)
table.columns[2].width = Inches(3)

style_table_cell(table.cell(0, 0), "任务", 16, WHITE, True, ACCENT_BLUE)
style_table_cell(table.cell(0, 1), "成功条件", 16, WHITE, True, ACCENT_BLUE)
style_table_cell(table.cell(0, 2), "判定方式", 16, WHITE, True, ACCENT_BLUE)
style_table_cell(table.cell(1, 0), "ObjectNav", 15, LIGHT_GRAY, False, CARD_BG)
style_table_cell(table.cell(1, 1), "目标物体在视野内且距离≤1.5m", 15, LIGHT_GRAY, False, CARD_BG)
style_table_cell(table.cell(1, 2), "可视性检测 + 距离计算", 15, LIGHT_GRAY, False, CARD_BG)
style_table_cell(table.cell(2, 0), "Pickup", 15, LIGHT_GRAY, False, CARD_BG)
style_table_cell(table.cell(2, 1), "目标物体被成功拿起", 15, LIGHT_GRAY, False, CARD_BG)
style_table_cell(table.cell(2, 2), "环境状态检测", 15, LIGHT_GRAY, False, CARD_BG)

# ============================================================
# 第4页：相关工作对比
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(8), Inches(0.8),
             "相关工作调研 (8篇代表性工作)", font_size=32, color=ACCENT_BLUE, bold=True)

table = add_table(slide, 9, 4, Inches(0.8), Inches(1.3), Inches(11.7), Inches(5.5))
table.columns[0].width = Inches(2)
table.columns[1].width = Inches(3)
table.columns[2].width = Inches(4.2)
table.columns[3].width = Inches(2.5)

# 表头
headers = ["类别", "论文/项目", "核心贡献", "与本项目关系"]
for j, h in enumerate(headers):
    style_table_cell(table.cell(0, j), h, 15, WHITE, True, ACCENT_BLUE)

# 数据
data = [
    ["仿真环境", "AI2-THOR (2017)", "4类室内场景,122个交互环境", "平台选型,支持导航+交互"],
    ["仿真环境", "Habitat-Sim (2019)", "高性能渲染,C++底层", "对比方案,交互能力弱"],
    ["VLA模型", "RT-2 (Google, 2023)", "VLM直接输出动作token", "架构设计核心灵感"],
    ["VLA模型", "RoboCat (2023)", "大规模预训练+微调", "SFT+PPO两阶段策略"],
    ["导航方法", "DD-PPO (2019)", "分布式PPO导航策略", "纯RL方法可行性证明"],
    ["导航方法", "CLIPort (2021)", "CLIP+空间注意力操作", "CLIP特征融合方式"],
    ["基础模型", "CLIP (2021) / SAM (2023)", "图文预训练/通用分割", "视觉编码组件"],
    ["基础模型", "Qwen3-VL 8B (2025)", "开源多模态大模型", "VLA决策基座模型"],
]

for i, row in enumerate(data):
    bg = CARD_BG if i % 2 == 0 else RGBColor(0x2D, 0x2D, 0x45)
    for j, val in enumerate(row):
        style_table_cell(table.cell(i+1, j), val, 13, LIGHT_GRAY, False, bg)

# ============================================================
# 第5页：系统架构图
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(8), Inches(0.8),
             "系统整体架构", font_size=32, color=ACCENT_BLUE, bold=True)

# 输入层
card1 = add_shape_bg(slide, Inches(0.5), Inches(1.2), Inches(3.5), Inches(1.2), CARD_BG)
add_text_box(slide, Inches(0.7), Inches(1.3), Inches(3), Inches(0.4),
             "输入层", font_size=18, color=ACCENT_GREEN, bold=True)
add_text_box(slide, Inches(0.7), Inches(1.7), Inches(3), Inches(0.6),
             "语言指令 + RGB图像 + Agent状态", font_size=14, color=LIGHT_GRAY)

# 感知层
card2 = add_shape_bg(slide, Inches(4.5), Inches(1.2), Inches(4), Inches(1.2), CARD_BG)
add_text_box(slide, Inches(4.7), Inches(1.3), Inches(3.5), Inches(0.4),
             "感知层", font_size=18, color=ACCENT_ORANGE, bold=True)
add_text_box(slide, Inches(4.7), Inches(1.7), Inches(3.5), Inches(0.6),
             "CLIP + MobileSAM + Position Encoder", font_size=14, color=LIGHT_GRAY)

# 融合层
card3 = add_shape_bg(slide, Inches(9), Inches(1.2), Inches(3.8), Inches(1.2), CARD_BG)
add_text_box(slide, Inches(9.2), Inches(1.3), Inches(3.5), Inches(0.4),
             "特征融合层", font_size=18, color=ACCENT_BLUE, bold=True)
add_text_box(slide, Inches(9.2), Inches(1.7), Inches(3.5), Inches(0.6),
             "投影对齐 → [seq, 4096]", font_size=14, color=LIGHT_GRAY)

# 箭头向下
arrow1 = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(6.2), Inches(2.5), Inches(0.8), Inches(0.6))
arrow1.fill.solid()
arrow1.fill.fore_color.rgb = ACCENT_BLUE
arrow1.line.fill.background()

# 决策层
card4 = add_shape_bg(slide, Inches(2), Inches(3.2), Inches(9.3), Inches(2), CARD_BG)
add_text_box(slide, Inches(2.2), Inches(3.3), Inches(4), Inches(0.5),
             "VLA决策层: Qwen3-VL 8B (QLoRA 4bit)", font_size=20, color=WHITE, bold=True)
add_bullet_list(slide, Inches(2.4), Inches(3.9), Inches(8.5), Inches(1.2),
                ["Multi-Modal Projector: visual tokens → LLM embedding空间",
                 "32层Transformer (LoRA adapter, rank=16, alpha=32)",
                 "Action Head: 6个特殊动作token加入词表(idx 32000-32005)"],
                font_size=15, color=LIGHT_GRAY, spacing=Pt(5))

arrow2 = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(6.2), Inches(5.3), Inches(0.8), Inches(0.6))
arrow2.fill.solid()
arrow2.fill.fore_color.rgb = ACCENT_BLUE
arrow2.line.fill.background()

# 输出层
card5 = add_shape_bg(slide, Inches(2), Inches(6), Inches(4), Inches(1), CARD_BG)
add_text_box(slide, Inches(2.2), Inches(6.1), Inches(3.5), Inches(0.4),
             "输出: 离散动作", font_size=18, color=ACCENT_GREEN, bold=True)
add_text_box(slide, Inches(2.2), Inches(6.5), Inches(3.5), Inches(0.5),
             "MoveAhead | Rotate | Pickup ...", font_size=14, color=LIGHT_GRAY)

card6 = add_shape_bg(slide, Inches(7), Inches(6), Inches(4.3), Inches(1), CARD_BG)
add_text_box(slide, Inches(7.2), Inches(6.1), Inches(3.8), Inches(0.4),
             "AI2-THOR环境", font_size=18, color=ACCENT_ORANGE, bold=True)
add_text_box(slide, Inches(7.2), Inches(6.5), Inches(3.8), Inches(0.5),
             "执行动作 → 返回新状态+reward", font_size=14, color=LIGHT_GRAY)

# ============================================================
# 第6页：关键技术 - 视觉编码
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(8), Inches(0.8),
             "关键技术: 多模态感知模块", font_size=32, color=ACCENT_BLUE, bold=True)

# CLIP
card1 = add_shape_bg(slide, Inches(0.8), Inches(1.3), Inches(5.5), Inches(2.5), CARD_BG)
add_text_box(slide, Inches(1), Inches(1.4), Inches(5), Inches(0.5),
             "CLIP ViT-B/32 - 全局场景理解", font_size=20, color=ACCENT_GREEN, bold=True)
add_bullet_list(slide, Inches(1.2), Inches(2), Inches(5), Inches(1.8),
                ["输入: 300×300 RGB图像", "输出: 512维全局特征向量", 
                 "作用: 捕获场景整体语义(房间类型、物体分布)", "参数量: ~150M, 推理速度快"],
                font_size=15, color=LIGHT_GRAY, spacing=Pt(5))

# MobileSAM
card2 = add_shape_bg(slide, Inches(7), Inches(1.3), Inches(5.5), Inches(2.5), CARD_BG)
add_text_box(slide, Inches(7.2), Inches(1.4), Inches(5), Inches(0.5),
             "MobileSAM - 目标检测与定位", font_size=20, color=ACCENT_ORANGE, bold=True)
add_bullet_list(slide, Inches(7.4), Inches(2), Inches(5), Inches(1.8),
                ["输入: 图像 + 目标类别提示(如'microwave')", 
                 "输出: mask + bounding box + 中心点坐标",
                 "作用: 提供'目标在哪'的空间定位信息", "轻量版~40MB, 适合实时推理"],
                font_size=15, color=LIGHT_GRAY, spacing=Pt(5))

# 位置编码
card3 = add_shape_bg(slide, Inches(0.8), Inches(4.2), Inches(11.7), Inches(2.8), CARD_BG)
add_text_box(slide, Inches(1), Inches(4.3), Inches(5), Inches(0.5),
             "Position Encoder - Agent状态编码", font_size=20, color=ACCENT_BLUE, bold=True)

# 流程图
add_text_box(slide, Inches(1.2), Inches(5), Inches(10), Inches(0.4),
             "[x, y, z, rotation]  →  归一化到[0,1]  →  Sinusoidal PE  →  Linear(256→4096)  →  [1, 4096]状态特征",
             font_size=16, color=LIGHT_GRAY)

add_bullet_list(slide, Inches(1.2), Inches(5.6), Inches(10), Inches(1.2),
                ["位置特征与视觉特征在输入LLM前进行拼接融合，形成完整的环境感知表示",
                 "类似Transformer中的Positional Encoding，将连续坐标映射到高维空间"],
                font_size=15, color=LIGHT_GRAY, spacing=Pt(5))

# ============================================================
# 第7页：模型设计 - Qwen3-VL
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(8), Inches(0.8),
             "模型设计: Qwen3-VL 8B + QLoRA", font_size=32, color=ACCENT_BLUE, bold=True)

# 左侧：模型参数
card1 = add_shape_bg(slide, Inches(0.8), Inches(1.3), Inches(5.5), Inches(3), CARD_BG)
add_text_box(slide, Inches(1), Inches(1.4), Inches(5), Inches(0.5),
             "基座模型参数", font_size=20, color=WHITE, bold=True)

table1 = add_table(slide, 6, 2, Inches(1.2), Inches(2.1), Inches(4.8), Inches(2))
table1.columns[0].width = Inches(2)
table1.columns[1].width = Inches(2.8)

params = [("基座模型", "Qwen3-VL 8B-Instruct"), ("参数量", "~8B"),
          ("视觉Encoder", "NaViT(动态分辨率)"), ("上下文窗口", "32K+"),
          ("FP16显存", "~16GB")]
for i, (k, v) in enumerate(params):
    bg = ACCENT_BLUE if i % 2 == 0 else CARD_BG
    style_table_cell(table1.cell(i, 0), k, 14, WHITE, True, bg)
    style_table_cell(table1.cell(i, 1), v, 14, LIGHT_GRAY, False, bg)

# 右侧：QLoRA配置
card2 = add_shape_bg(slide, Inches(7), Inches(1.3), Inches(5.5), Inches(3), CARD_BG)
add_text_box(slide, Inches(7.2), Inches(1.4), Inches(5), Inches(0.5),
             "QLoRA 微调配置 (4bit)", font_size=20, color=ACCENT_GREEN, bold=True)

table2 = add_table(slide, 6, 2, Inches(7.4), Inches(2.1), Inches(4.8), Inches(2))
table2.columns[0].width = Inches(2)
table2.columns[1].width = Inches(2.8)

lora = [("量化方式", "NF4 4bit"), ("LoRA rank", "16"),
        ("LoRA alpha", "32"), ("可训练参数", "~7M (0.1%)"),
        ("训练显存", "~16-18GB")]
for i, (k, v) in enumerate(lora):
    bg = ACCENT_GREEN if i % 2 == 0 else CARD_BG
    style_table_cell(table2.cell(i, 0), k, 14, WHITE, True, bg)
    style_table_cell(table2.cell(i, 1), v, 14, LIGHT_GRAY, False, bg)

# 动作token
card3 = add_shape_bg(slide, Inches(0.8), Inches(4.7), Inches(11.7), Inches(2.3), CARD_BG)
add_text_box(slide, Inches(1), Inches(4.8), Inches(5), Inches(0.5),
             "动作Token化 (6个离散动作)", font_size=20, color=ACCENT_ORANGE, bold=True)

tokens = ["<ACT_MoveAhead> (idx=32000)", "<ACT_RotateLeft> (idx=32001)", 
          "<ACT_RotateRight> (idx=32002)", "<ACT_LookUp> (idx=32003)",
          "<ACT_LookDown> (idx=32004)", "<ACT_Pickup> (idx=32005)"]
for i, token in enumerate(tokens):
    x = Inches(1.2) + (i % 3) * Inches(3.8)
    y = Inches(5.5) + (i // 3) * Inches(0.5)
    add_text_box(slide, x, y, Inches(3.5), Inches(0.5),
                 token, font_size=15, color=LIGHT_GRAY)

# ============================================================
# 第8页：训练策略
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(8), Inches(0.8),
             "训练策略: SFT → PPO 两阶段", font_size=32, color=ACCENT_BLUE, bold=True)

# 阶段1 SFT
card1 = add_shape_bg(slide, Inches(0.8), Inches(1.3), Inches(5.5), Inches(3.2), CARD_BG)
add_text_box(slide, Inches(1), Inches(1.4), Inches(5), Inches(0.5),
             "阶段1: 监督微调 (SFT)", font_size=22, color=ACCENT_GREEN, bold=True)
add_bullet_list(slide, Inches(1.2), Inches(2), Inches(5), Inches(2.5),
                ["数据: 启发式策略(A*寻路)生成专家演示", "规模: ~5000条轨迹 × 15步 ≈ 75K样本",
                 "损失: 交叉熵(预测动作 vs 专家动作)", "训练: 3 epochs, 学习率 2e-4",
                 "目标: 让模型学会'基本导航能力'"],
                font_size=15, color=LIGHT_GRAY, spacing=Pt(5))

# 阶段2 PPO
card2 = add_shape_bg(slide, Inches(7), Inches(1.3), Inches(5.5), Inches(3.2), CARD_BG)
add_text_box(slide, Inches(7.2), Inches(1.4), Inches(5), Inches(0.5),
             "阶段2: PPO强化学习", font_size=22, color=ACCENT_ORANGE, bold=True)
add_bullet_list(slide, Inches(7.4), Inches(2), Inches(5), Inches(2.5),
                ["算法: PPO-Clip (稳定版)", "折扣因子 γ=0.99, GAE λ=0.95",
                 "Clip范围 ε=0.2", "在SFT基础上优化长期回报",
                 "目标: 提升成功率 + 路径效率"],
                font_size=15, color=LIGHT_GRAY, spacing=Pt(5))

# 奖励函数
card3 = add_shape_bg(slide, Inches(0.8), Inches(4.8), Inches(11.7), Inches(2.2), CARD_BG)
add_text_box(slide, Inches(1), Inches(4.9), Inches(5), Inches(0.5),
             "奖励函数设计", font_size=20, color=WHITE, bold=True)

table = add_table(slide, 2, 6, Inches(1), Inches(5.5), Inches(11.3), Inches(1.2))
table.columns[0].width = Inches(1.8)
table.columns[1].width = Inches(1.8)
table.columns[2].width = Inches(1.8)
table.columns[3].width = Inches(2)
table.columns[4].width = Inches(2)
table.columns[5].width = Inches(1.9)

rewards = ["到达目标\n+10.0", "成功拾取\n+5.0", "每走一步\n-0.1", "靠近目标\n+0.05/步", "远离目标\n-0.05/步", "碰撞/越界\n-1.0"]
for j, r in enumerate(rewards):
    bg = ACCENT_BLUE if j < 2 else (RGBColor(0xCC, 0x44, 0x00) if j == 5 else CARD_BG)
    style_table_cell(table.cell(0, j), r, 13, WHITE, True, bg)

# ============================================================
# 第9页：评测方案
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(5), Inches(0.8),
             "评测指标与消融实验", font_size=32, color=ACCENT_BLUE, bold=True)

# 主要指标
card1 = add_shape_bg(slide, Inches(0.8), Inches(1.3), Inches(5.5), Inches(3), CARD_BG)
add_text_box(slide, Inches(1), Inches(1.4), Inches(5), Inches(0.5),
             "主要指标", font_size=20, color=WHITE, bold=True)

add_text_box(slide, Inches(1.2), Inches(2), Inches(5), Inches(0.5),
             "1. 成功率 (Success Rate)", font_size=18, color=ACCENT_GREEN, bold=True)
add_text_box(slide, Inches(1.4), Inches(2.5), Inches(4.5), Inches(0.5),
             "SR = N_success / N_total", font_size=16, color=LIGHT_GRAY)

add_text_box(slide, Inches(1.2), Inches(3), Inches(5), Inches(0.5),
             "2. SPL (Success weighted by Path Length)", font_size=18, color=ACCENT_GREEN, bold=True)
add_text_box(slide, Inches(1.4), Inches(3.5), Inches(4.5), Inches(0.7),
             "SPL = (1/N) Σ Sᵢ · lᵢ*/max(lᵢ*, lᵢ)\n同时衡量成功率和路径效率", font_size=15, color=LIGHT_GRAY)

# 消融实验
card2 = add_shape_bg(slide, Inches(7), Inches(1.3), Inches(5.5), Inches(3), CARD_BG)
add_text_box(slide, Inches(7.2), Inches(1.4), Inches(5), Inches(0.5),
             "消融实验设计", font_size=20, color=ACCENT_ORANGE, bold=True)

table = add_table(slide, 4, 4, Inches(7.2), Inches(2), Inches(5.1), Inches(2))
table.columns[0].width = Inches(1.3)
table.columns[1].width = Inches(1)
table.columns[2].width = Inches(1)
table.columns[3].width = Inches(1.8)

style_table_cell(table.cell(0, 0), "实验组", 13, WHITE, True, ACCENT_BLUE)
style_table_cell(table.cell(0, 1), "视觉", 13, WHITE, True, ACCENT_BLUE)
style_table_cell(table.cell(0, 2), "语言", 13, WHITE, True, ACCENT_BLUE)
style_table_cell(table.cell(0, 3), "对比目的", 13, WHITE, True, ACCENT_BLUE)

style_table_cell(table.cell(1, 0), "Baseline-A", 13, LIGHT_GRAY, False, CARD_BG)
style_table_cell(table.cell(1, 1), "CLIP", 13, LIGHT_GRAY, False, CARD_BG)
style_table_cell(table.cell(1, 2), "✗", 13, LIGHT_GRAY, False, CARD_BG)
style_table_cell(table.cell(1, 3), "纯视觉导航", 13, LIGHT_GRAY, False, CARD_BG)

style_table_cell(table.cell(2, 0), "Baseline-B", 13, LIGHT_GRAY, False, CARD_BG)
style_table_cell(table.cell(2, 1), "CLIP+SAM", 13, LIGHT_GRAY, False, CARD_BG)
style_table_cell(table.cell(2, 2), "✓", 13, LIGHT_GRAY, False, CARD_BG)
style_table_cell(table.cell(2, 3), "无位置编码", 13, LIGHT_GRAY, False, CARD_BG)

style_table_cell(table.cell(3, 0), "Ours", 13, WHITE, True, ACCENT_GREEN)
style_table_cell(table.cell(3, 1), "CLIP+SAM", 13, WHITE, True, ACCENT_GREEN)
style_table_cell(table.cell(3, 2), "✓+位置", 13, WHITE, True, ACCENT_GREEN)
style_table_cell(table.cell(3, 3), "完整模型", 13, WHITE, True, ACCENT_GREEN)

# 次要指标
card3 = add_shape_bg(slide, Inches(0.8), Inches(4.6), Inches(11.7), Inches(2.4), CARD_BG)
add_text_box(slide, Inches(1), Inches(4.7), Inches(5), Inches(0.5),
             "次要指标", font_size=20, color=WHITE, bold=True)

table2 = add_table(slide, 2, 4, Inches(1), Inches(5.3), Inches(11.3), Inches(1.3))
table2.columns[0].width = Inches(2.5)
table2.columns[1].width = Inches(3)
table2.columns[2].width = Inches(2.5)
table2.columns[3].width = Inches(3.3)

style_table_cell(table2.cell(0, 0), "平均步数", 14, WHITE, True, ACCENT_BLUE)
style_table_cell(table2.cell(0, 1), "成功episode的平均动作数", 14, LIGHT_GRAY, False, CARD_BG)
style_table_cell(table2.cell(0, 2), "路径效率", 14, WHITE, True, ACCENT_BLUE)
style_table_cell(table2.cell(0, 3), "实际路径/最优路径比值", 14, LIGHT_GRAY, False, CARD_BG)

# ============================================================
# 第10页：实验计划
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(5), Inches(0.8),
             "实验计划: 算力与时间", font_size=32, color=ACCENT_BLUE, bold=True)

# 算力
card1 = add_shape_bg(slide, Inches(0.8), Inches(1.3), Inches(5.5), Inches(2.5), CARD_BG)
add_text_box(slide, Inches(1), Inches(1.4), Inches(5), Inches(0.5),
             "算力预算", font_size=20, color=WHITE, bold=True)

table = add_table(slide, 5, 3, Inches(1), Inches(2.1), Inches(5), Inches(1.6))
table.columns[0].width = Inches(1.5)
table.columns[1].width = Inches(2)
table.columns[2].width = Inches(1.5)

style_table_cell(table.cell(0, 0), "组件", 14, WHITE, True, ACCENT_BLUE)
style_table_cell(table.cell(0, 1), "规格", 14, WHITE, True, ACCENT_BLUE)
style_table_cell(table.cell(0, 2), "用途", 14, WHITE, True, ACCENT_BLUE)

specs = [("GPU", "RTX 4090 (24GB)", "模型训练"), ("CPU", "8核以上", "环境仿真"),
         ("内存", "32GB", "数据加载"), ("存储", "50GB", "模型+数据")]
for i, (a, b, c) in enumerate(specs):
    bg = CARD_BG
    style_table_cell(table.cell(i+1, 0), a, 13, LIGHT_GRAY, False, bg)
    style_table_cell(table.cell(i+1, 1), b, 13, LIGHT_GRAY, False, bg)
    style_table_cell(table.cell(i+1, 2), c, 13, LIGHT_GRAY, False, bg)

# 训练时间
card2 = add_shape_bg(slide, Inches(7), Inches(1.3), Inches(5.5), Inches(2.5), CARD_BG)
add_text_box(slide, Inches(7.2), Inches(1.4), Inches(5), Inches(0.5),
             "训练时间估算", font_size=20, color=WHITE, bold=True)

table2 = add_table(slide, 4, 3, Inches(7.2), Inches(2.1), Inches(5), Inches(1.4))
table2.columns[0].width = Inches(1.5)
table2.columns[1].width = Inches(1.8)
table2.columns[2].width = Inches(1.7)

style_table_cell(table2.cell(0, 0), "阶段", 14, WHITE, True, ACCENT_BLUE)
style_table_cell(table2.cell(0, 1), "步数", 14, WHITE, True, ACCENT_BLUE)
style_table_cell(table2.cell(0, 2), "预估时间", 14, WHITE, True, ACCENT_BLUE)

times = [("SFT", "75K×3 epochs", "~8小时"), ("PPO", "50K steps", "~10小时"),
         ("评估", "300 episodes", "~5分钟")]
for i, (a, b, c) in enumerate(times):
    style_table_cell(table2.cell(i+1, 0), a, 13, LIGHT_GRAY, False, CARD_BG)
    style_table_cell(table2.cell(i+1, 1), b, 13, LIGHT_GRAY, False, CARD_BG)
    style_table_cell(table2.cell(i+1, 2), c, 13, LIGHT_GRAY, False, CARD_BG)

# 时间线
card3 = add_shape_bg(slide, Inches(0.8), Inches(4.2), Inches(11.7), Inches(2.8), CARD_BG)
add_text_box(slide, Inches(1), Inches(4.3), Inches(5), Inches(0.5),
             "时间计划 (第9-18周)", font_size=20, color=WHITE, bold=True)

timeline = [("第9-10周", "开题报告 + 展示", ACCENT_BLUE),
            ("第11-12周", "环境搭建 + 数据生成", ACCENT_GREEN),
            ("第13-14周", "模型实现 + SFT训练", ACCENT_ORANGE),
            ("第15周", "PPO训练 + 调优", ACCENT_BLUE),
            ("第16-17周", "评估 + 可视化 + 消融实验", ACCENT_GREEN),
            ("第18周", "最终报告 + 期末展示", ACCENT_ORANGE)]

for i, (week, task, color) in enumerate(timeline):
    y = Inches(4.9) + i * Inches(0.35)
    add_text_box(slide, Inches(1.2), y, Inches(2), Inches(0.35),
                 week, font_size=14, color=color, bold=True)
    add_text_box(slide, Inches(3.5), y, Inches(8), Inches(0.35),
                 task, font_size=14, color=LIGHT_GRAY)

# ============================================================
# 第11页：预期成果
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(5), Inches(0.8),
             "预期成果", font_size=32, color=ACCENT_BLUE, bold=True)

# 量化指标
card1 = add_shape_bg(slide, Inches(0.8), Inches(1.3), Inches(5.5), Inches(2.5), CARD_BG)
add_text_box(slide, Inches(1), Inches(1.4), Inches(5), Inches(0.5),
             "量化目标", font_size=20, color=WHITE, bold=True)

table = add_table(slide, 3, 3, Inches(1), Inches(2.1), Inches(5), Inches(1.3))
table.columns[0].width = Inches(1.8)
table.columns[1].width = Inches(1.5)
table.columns[2].width = Inches(1.7)

style_table_cell(table.cell(0, 0), "指标", 14, WHITE, True, ACCENT_BLUE)
style_table_cell(table.cell(0, 1), "Seen场景", 14, WHITE, True, ACCENT_BLUE)
style_table_cell(table.cell(0, 2), "Unseen场景", 14, WHITE, True, ACCENT_BLUE)

style_table_cell(table.cell(1, 0), "成功率", 15, LIGHT_GRAY, False, CARD_BG)
style_table_cell(table.cell(1, 1), "> 50%", 15, ACCENT_GREEN, True, CARD_BG)
style_table_cell(table.cell(1, 2), "> 30%", 15, ACCENT_ORANGE, True, CARD_BG)

style_table_cell(table.cell(2, 0), "SPL", 15, LIGHT_GRAY, False, CARD_BG)
style_table_cell(table.cell(2, 1), "> 0.3", 15, ACCENT_GREEN, True, CARD_BG)
style_table_cell(table.cell(2, 2), "> 0.2", 15, ACCENT_ORANGE, True, CARD_BG)

# 交付物
card2 = add_shape_bg(slide, Inches(7), Inches(1.3), Inches(5.5), Inches(2.5), CARD_BG)
add_text_box(slide, Inches(7.2), Inches(1.4), Inches(5), Inches(0.5),
             "交付物", font_size=20, color=ACCENT_GREEN, bold=True)
add_bullet_list(slide, Inches(7.4), Inches(2), Inches(5), Inches(1.8),
                ["✅ 完整代码仓库(GitHub)", "✅ 开题报告(5页) + 最终报告(10-15页)",
                 "✅ 10+ episode轨迹可视化", "✅ 失败案例分类分析",
                 "✅ 3-5分钟演示视频"],
                font_size=15, color=LIGHT_GRAY, spacing=Pt(4))

# 风险应对
card3 = add_shape_bg(slide, Inches(0.8), Inches(4.2), Inches(11.7), Inches(2.8), CARD_BG)
add_text_box(slide, Inches(1), Inches(4.3), Inches(5), Inches(0.5),
             "风险与应对", font_size=20, color=ACCENT_ORANGE, bold=True)

table2 = add_table(slide, 4, 3, Inches(1), Inches(5), Inches(11.3), Inches(1.6))
table2.columns[0].width = Inches(3.5)
table2.columns[1].width = Inches(1.5)
table2.columns[2].width = Inches(6.3)

style_table_cell(table2.cell(0, 0), "风险", 14, WHITE, True, ACCENT_BLUE)
style_table_cell(table2.cell(0, 1), "概率", 14, WHITE, True, ACCENT_BLUE)
style_table_cell(table2.cell(0, 2), "应对方案", 14, WHITE, True, ACCENT_BLUE)

risks = [("AI2-THOR在WSL运行不稳定", "中", "使用headless模式; 备用Docker方案"),
         ("PPO训练不收敛", "中", "增加SFT数据; 调整奖励; 使用GAE"),
         ("训练时间超出预算", "中", "减少场景数; 使用云平台(AutoDL)")]
for i, (a, b, c) in enumerate(risks):
    style_table_cell(table2.cell(i+1, 0), a, 13, LIGHT_GRAY, False, CARD_BG)
    style_table_cell(table2.cell(i+1, 1), b, 13, ACCENT_ORANGE, True, CARD_BG)
    style_table_cell(table2.cell(i+1, 2), c, 13, LIGHT_GRAY, False, CARD_BG)

# ============================================================
# 第12页：Q&A
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

line1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2), Inches(3), Inches(9.3), Inches(0.03))
line1.fill.solid()
line1.fill.fore_color.rgb = ACCENT_BLUE
line1.line.fill.background()

add_text_box(slide, Inches(2), Inches(1.5), Inches(9.3), Inches(1.5),
             "Q & A", font_size=54, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

line2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2), Inches(4.2), Inches(9.3), Inches(0.03))
line2.fill.solid()
line2.fill.fore_color.rgb = ACCENT_BLUE
line2.line.fill.background()

add_text_box(slide, Inches(2), Inches(4.5), Inches(9.3), Inches(1),
             "感谢聆听！", font_size=28, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(2), Inches(5.8), Inches(9.3), Inches(1),
             "GitHub: https://github.com/your-org/mllm-agent", font_size=16, 
             color=MEDIUM_GRAY, alignment=PP_ALIGN.CENTER)

# ========== 保存 ==========
output_path = r"d:\experiment\mLLM\docs\开题展示.pptx"
prs.save(output_path)
print(f"PPT已保存至: {output_path}")
